import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import Patch
import numpy as np
import io
import seaborn as sns
import streamlit.components.v1 as components
from PIL import Image
from docx import Document
from docx.shared import Inches
import tempfile
# --- FIX: Tắt giới hạn pixel để tránh lỗi DecompressionBombError ---
Image.MAX_IMAGE_PIXELS = None

# --- PAGE CONFIG ---
st.set_page_config(page_title="Quality & Scrap Dashboard", layout="wide")
st.title("📊 Production Quality Yield & Tail Scrap Analysis")
st.markdown("---")

# --- FIX: Tối ưu DPI cho Web để tiết kiệm RAM (vẫn giữ 300 DPI cho tải xuống) ---
plt.rcParams['figure.dpi'] = 120
plt.rcParams['savefig.dpi'] = 300
plt.rcParams['savefig.bbox'] = 'tight'

# --- SIDEBAR: DYNAMIC SPEC LIMITS INPUT PER THICKNESS ---
st.sidebar.header("⚙️ Spec Limits (Control)")
st.sidebar.info("Limits apply from Q4 2025 onwards. Configured per Thickness.")

GLOBAL_SPECS = {0.5: {}, 0.6: {}, 0.8: {}}
target_thicks = [0.5, 0.6, 0.8]

for t in target_thicks:
    with st.sidebar.expander(f"📏 Limits for {t}mm"):
        for feat in ['YS', 'TS', 'EL', 'YPE']:
            st.caption(f"**{feat}**")
            c1, c2, c3 = st.columns(3)
            min_v = c1.number_input("Min", value=None, key=f"{t}_{feat}_min", label_visibility="collapsed", placeholder="Min")
            max_v = c2.number_input("Max", value=None, key=f"{t}_{feat}_max", label_visibility="collapsed", placeholder="Max")
            tgt_v = c3.number_input("Tgt", value=None, key=f"{t}_{feat}_tgt", label_visibility="collapsed", placeholder="Tgt")
            GLOBAL_SPECS[t][feat] = {'min': min_v, 'max': max_v, 'target': tgt_v}

uploaded_file = st.file_uploader("Upload Production Data (.xlsx)", type=["xlsx"])

if uploaded_file is not None:
    df = pd.read_excel(uploaded_file)
    df.columns = df.columns.astype(str).str.strip()

    # --- 1. DATA PRE-PROCESSING ---
    # Handle Dates & Time Grouping First
    date_key = '烤三生產日期' 
    if date_key in df.columns:
        d_str = df[date_key].astype(str).str.replace(r'\.0$', '', regex=True).str.strip()
        df['Production_Date'] = pd.to_datetime(d_str, format='%Y%m%d', errors='coerce')
        
        def categorize_period(d):
            if pd.isnull(d): return "Unknown"
            y = d.year
            q3_s, q3_e = pd.Timestamp(2025, 6, 29), pd.Timestamp(2025, 9, 30)
            
            if y == 2024: return "2024 (Full Year)"
            if y == 2025:
                if d < q3_s: return "2025 H1 (Until 06/28)"
                if q3_s <= d <= q3_e: return "2025 Q3 (06/29 - 09/30)"
                return d.strftime('%Y-%m')
            if y >= 2026: return d.strftime('%Y-%m')
            return "Other"
            
        df['Time_Group'] = df['Production_Date'].apply(categorize_period)
        df = df[df['Time_Group'] != "Other"]
        
        df_25 = df[df['Production_Date'].dt.year == 2025].copy()
        if not df_25.empty:
            df_25['Time_Group'] = "2025 (Full Year)"
            df = pd.concat([df, df_25], ignore_index=True)
    else:
        df['Time_Group'] = "Unknown"

    # Robust Quality Grade Mapping
    base_grades = ['A-B+', 'A-B', 'A-B-', 'B+', 'B']
    for g in base_grades:
        match_cols = []
        for c in df.columns:
            c_str = str(c).strip()
            if c_str == g or c_str == f"{g}個數" or c_str.startswith(f"{g}."):
                match_cols.append(c)
        df[g] = df[match_cols].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1) if match_cols else 0

    df['Total_Qty'] = df[base_grades].sum(axis=1)
    df['Severe_Bad_Qty'] = df[['B+', 'B']].sum(axis=1)
    df['Acceptable_Qty'] = df['Total_Qty'] - df['Severe_Bad_Qty']
    
    target_grades = ['A-B+', 'A-B']
    df['Valid_Qty'] = df[target_grades].sum(axis=1)

    # STORE ORIGINAL DATA FOR GRADE DISTRIBUTION (Unfiltered by thickness)
    df_global_grades = df[df['Total_Qty'] > 0].copy()

    # Handle Thickness & Filtering for Detailed Charts
    if 'Actual_Thickness' not in df.columns:
        if 'Thickness' in df.columns:
            df.rename(columns={'Thickness': 'Actual_Thickness'}, inplace=True)
        elif '厚度' in df.columns:
            df.rename(columns={'厚度': 'Actual_Thickness'}, inplace=True)
        else:
            for i, c in enumerate(df.columns):
                if '型式' in c and i > 0:
                    df.rename(columns={df.columns[i - 1]: 'Actual_Thickness'}, inplace=True)
                    break
    
    if 'Actual_Thickness' in df.columns:
        df['Actual_Thickness'] = pd.to_numeric(df['Actual_Thickness'], errors='coerce')
        
        def map_thickness(val):
            if pd.isna(val): return None
            v = round(float(val), 2)
            if v in [0.47, 0.50]: return 0.5
            if v in [0.53, 0.54, 0.57, 0.58, 0.60]: return 0.6
            if v in [0.63, 0.75, 0.76, 0.77, 0.80]: return 0.8
            return None 
            
        df['Standard_Thickness'] = df['Actual_Thickness'].apply(map_thickness)
        df = df.dropna(subset=['Standard_Thickness'])
        df['Actual_Thickness'] = df['Standard_Thickness']
        df = df.drop(columns=['Standard_Thickness'])
    else:
        df['Actual_Thickness'] = 0.0
        df = df.iloc[0:0]

    if '熱軋材質' in df.columns:
        df['HR_Material'] = df['熱軋材質'].astype(str).str.strip().replace(['nan', ''], 'Unknown')
    else:
        df['HR_Material'] = 'Unknown'

    df.rename(columns={'烤漆降伏強度': 'YS', '烤漆抗拉強度': 'TS', '伸長率': 'EL'}, inplace=True)
    mech_features = ['YS', 'TS', 'EL', 'YPE', 'HARDNESS']
    for f in mech_features:
        if f in df.columns:
            df[f] = pd.to_numeric(df[f], errors='coerce')

    LEN_COL = '實測長度'
    SCRAP_COL = '尾料剔退'
    if LEN_COL in df.columns:
        df[LEN_COL] = pd.to_numeric(df[LEN_COL], errors='coerce').fillna(0)
    if SCRAP_COL in df.columns:
        df[SCRAP_COL] = pd.to_numeric(df[SCRAP_COL], errors='coerce').fillna(0)

    df = df[(df['Total_Qty'] > 0) | (df.get(LEN_COL, 0) > 0) | (df.get(SCRAP_COL, 0) > 0)]

    sns.set_theme(style="whitegrid")
    solid_colors = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd', '#8c564b']

    def get_sort_key(x):
        if "2024 (Full Year)" in x: return "2024-00"
        if "2025 H1" in x: return "2025-00a"
        if "2025 Q3" in x: return "2025-00b"
        if "2025 (Full Year)" in x: return "2025-99" 
        return x

    def add_chart_border(ax):
        for spine in ax.spines.values():
            spine.set_visible(True)
            spine.set_color('#333333')
            spine.set_linewidth(1.0)

    # --- SPC & CAPABILITY HELPER FUNCTIONS ---
    def is_valid_for_control(period_label):
        if "2024" in period_label or "2025 H1" in period_label or "2025 Q3" in period_label or "2025 (Full Year)" in period_label:
            return False
        return True

    def calc_capability(values, feat, period_label, thickness):
        vals = np.array(values, dtype=float)
        vals = vals[~np.isnan(vals)]
        
        if thickness == 'Overall':
            return {'mean': np.mean(vals) if len(vals) > 0 else 0, 'std': np.std(vals, ddof=1) if len(vals) > 1 else 0, 'n': len(vals),
                    'Cp': None, 'Cpk': None, 'Ca': None, 'LSL': None, 'USL': None, 'Target': None}

        if len(vals) < 2: return None
        mu  = np.mean(vals)
        std = np.std(vals, ddof=1)
        if std == 0: return None

        spec = GLOBAL_SPECS.get(thickness, {}).get(feat, {})
        lsl  = spec.get('min')
        usl  = spec.get('max')
        tgt  = spec.get('target')

        result = {'mean': mu, 'std': std, 'n': len(vals),
                  'Cp': None, 'Cpk': None, 'Ca': None,
                  'LSL': lsl, 'USL': usl, 'Target': tgt}

        if not is_valid_for_control(period_label):
            return result

        if lsl is not None and usl is not None:
            cp   = (usl - lsl) / (6 * std)
            cpu  = (usl - mu)  / (3 * std)
            cpl  = (mu  - lsl) / (3 * std)
            cpk  = min(cpu, cpl)
            result['Cp']  = round(cp,  3)
            result['Cpk'] = round(cpk, 3)
            if tgt is not None:
                ca = (mu - tgt) / ((usl - lsl) / 2) * 100
                result['Ca'] = round(ca, 2)
            else:
                mid = (usl + lsl) / 2
                ca  = (mu - mid) / ((usl - lsl) / 2) * 100
                result['Ca'] = round(ca, 2)
        elif lsl is not None:
            cpl = (mu - lsl) / (3 * std)
            result['Cp']  = round(cpl, 3)
            result['Cpk'] = round(cpl, 3)
        elif usl is not None:
            cpu = (usl - mu) / (3 * std)
            result['Cp']  = round(cpu, 3)
            result['Cpk'] = round(cpu, 3)

        return result

    def cpk_color(cpk):
        if cpk is None: return '#888888'
        if cpk >= 1.67: return '#2e7d32'   
        if cpk >= 1.33: return '#66bb6a'   
        if cpk >= 1.00: return '#ffa726'   
        return '#d62728'                    

    def cpk_label(cpk, period_label, thickness):
        if thickness == 'Overall': return 'Mixed Specs'
        if not is_valid_for_control(period_label): return 'Not Monitored'
        if cpk is None: return 'N/A'
        if cpk >= 1.67: return '✅ Excellent'
        if cpk >= 1.33: return '✅ Capable'
        if cpk >= 1.00: return '⚠️ Marginal'
        return '❌ Not Capable'

    def render_capability_badge(cap, feat, period_label, thickness):
        if cap is None: return
        
        mu_v   = f"{cap['mean']:.2f}"
        std_v  = f"{cap['std']:.3f}"
        
        if thickness == 'Overall':
            cp_v, cpk_v, ca_v = 'N/A', 'N/A', 'N/A'
            clr, lbl = '#888888', 'Mixed Thickness (No Global Limit)'
            lsl_v, usl_v = '—', '—'
        elif is_valid_for_control(period_label):
            cp_v   = f"{cap['Cp']:.3f}"   if cap['Cp']  is not None else 'N/A'
            cpk_v  = f"{cap['Cpk']:.3f}"  if cap['Cpk'] is not None else 'N/A'
            ca_v   = f"{cap['Ca']:.1f}%"  if cap['Ca']  is not None else 'N/A'
            clr    = cpk_color(cap['Cpk'])
            lbl    = cpk_label(cap['Cpk'], period_label, thickness)
            lsl_v  = str(cap['LSL']) if cap['LSL'] is not None else '—'
            usl_v  = str(cap['USL']) if cap['USL'] is not None else '—'
        else:
            cp_v, cpk_v, ca_v = 'N/A', 'N/A', 'N/A'
            clr, lbl = '#888888', 'Pre-Q4 2025 (Not Monitored)'
            lsl_v, usl_v = '—', '—'

        html_badge = f"""
        <div style="background:#f8f9fa;border-left:5px solid {clr};
                    border-radius:6px;padding:8px 14px;margin:4px 0 10px 0;
                    font-family:monospace;font-size:13px;line-height:1.8;">
          <span style="font-size:14px;font-weight:bold;color:{clr};">{lbl}</span>
          &nbsp;&nbsp;|&nbsp;&nbsp;
          <b>LSL</b>: {lsl_v} &nbsp; <b>USL</b>: {usl_v}
          &nbsp;&nbsp;|&nbsp;&nbsp;
          <b>n</b>: {cap['n']} &nbsp;
          <b>Mean</b>: {mu_v} &nbsp;
          <b>Std</b>: {std_v}
          <br>
          <b style="color:{clr};">Cpk = {cpk_v}</b> &nbsp;&nbsp;
          <b>Cp = {cp_v}</b> &nbsp;&nbsp;
          <b>Ca = {ca_v}</b>
        </div>
        """
        st.markdown(html_badge, unsafe_allow_html=True)

    def build_capability_summary(df_src, feat, label, thick):
        vals = df_src[feat].dropna().values if feat in df_src.columns else []
        cap = calc_capability(vals, feat, label, thick)
        if cap is None or cap['n'] < 2: return None
        
        if not is_valid_for_control(label):
            return {
                'Period': label, 'Thickness': thick, 'Feature': feat, 'n': cap['n'],
                'Mean': round(cap['mean'], 2), 'Std': round(cap['std'], 3),
                'LSL': None, 'USL': None, 'Ca (%)': None, 'Cp': None, 'Cpk': None,
                'Verdict': 'Not Monitored'
            }

        return {
            'Period': label, 'Thickness': thick, 'Feature': feat, 'n': cap['n'],
            'Mean': round(cap['mean'], 2), 'Std': round(cap['std'], 3),
            'LSL': cap['LSL'], 'USL': cap['USL'], 'Ca (%)': cap['Ca'],
            'Cp': cap['Cp'], 'Cpk': cap['Cpk'],
            'Verdict': cpk_label(cap['Cpk'], label, thick).replace('✅ ', '').replace('⚠️ ', '').replace('❌ ', '')
        }

    global_x_bounds = {}
    for feat in ['YS', 'TS', 'EL', 'YPE']:
        if feat in df.columns:
            vd = df[[feat, 'Valid_Qty']].dropna().copy()
            vd = vd[vd['Valid_Qty'] > 0]
            if not vd.empty:
                q1, q99 = np.percentile(vd[feat], 1), np.percentile(vd[feat], 99)
                global_x_bounds[feat] = (q1 - (q99 - q1) * 0.25, q99 + (q99 - q1) * 0.25)

    def get_shared_y(data, features):
        max_y = 0
        for f in features:
            if f in data.columns:
                vd = data.dropna(subset=[f])
                if not vd.empty:
                    cnts, _ = np.histogram(vd[f], bins=15, weights=vd['Total_Qty'])
                    max_y = max(max_y, cnts.max())
        return max_y * 1.35 if max_y > 0 else 50

    def plot_dist(ax, data, feat, title, y_lim, period_label, thickness):
        c_map = {'A-B+': '#2ca02c', 'A-B': '#1f77b4', 'A-B-': '#ff7f0e', 'B+': '#9467bd', 'B': '#d62728'}
        fmin, fmax = global_x_bounds.get(feat, (data[feat].min() if not data.empty else 0, data[feat].max() if not data.empty else 100))
        if fmin == fmax: fmax += 1
        
        v_l, w_l, clrs, m_info = [], [], [], []
        for g in base_grades:
            if g in data.columns:
                td = data[[feat, g]].dropna()
                td = td[td[g] > 0]
                if not td.empty:
                    v_l.append(td[feat].values)
                    w_l.append(td[g].values)
                    clrs.append(c_map[g])
                    m = np.average(td[feat].values, weights=td[g].values)
                    ax.axvline(m, color=c_map[g], ls='--', lw=1.2)
                    m_info.append({'v': m, 'c': c_map[g], 'label': g})

        if v_l:
            ax.hist(v_l, bins=np.linspace(fmin, fmax, 16), weights=w_l, color=clrs, stacked=True, edgecolor='white', alpha=0.7)
            m_info.sort(key=lambda x: x['v'])
            x_range = fmax - fmin
            min_gap = x_range * 0.045
            positions = [info['v'] for info in m_info]
            for _ in range(50):
                moved = False
                for i in range(1, len(positions)):
                    if positions[i] - positions[i - 1] < min_gap:
                        mid = (positions[i] + positions[i - 1]) / 2
                        positions[i - 1] = mid - min_gap / 2
                        positions[i] = mid + min_gap / 2
                        moved = True
                if not moved: break
            
            y_levels = [y_lim * (0.92 - (i % 4) * 0.13) for i in range(len(m_info))]
            for i, info in enumerate(m_info):
                x_pos = positions[i]
                y_pos = y_levels[i]
                ax.annotate(
                    f"{info['v']:.1f}",
                    xy=(info['v'], y_pos * 0.6), xytext=(x_pos, y_pos),
                    color='white', fontweight='bold', fontsize=8, ha='center', va='center',
                    bbox=dict(facecolor=info['c'], alpha=0.85, boxstyle='round,pad=0.25'),
                    arrowprops=dict(arrowstyle='-', color=info['c'], lw=1.0, alpha=0.6) if abs(x_pos - info['v']) > min_gap * 0.3 else None
                )

        if thickness != 'Overall' and is_valid_for_control(period_label):
            spec = GLOBAL_SPECS.get(thickness, {}).get(feat, {})
            lsl, usl, tgt = spec.get('min'), spec.get('max'), spec.get('target')
            y_top = y_lim * 0.98
            if lsl is not None:
                ax.axvline(lsl, color='darkred', lw=2, ls='-', zorder=3)
                ax.text(lsl, y_top, f' LSL\n {lsl}', color='darkred', fontsize=7.5, fontweight='bold', va='top', ha='left')
            if usl is not None:
                ax.axvline(usl, color='darkred', lw=2, ls='-', zorder=3)
                ax.text(usl, y_top, f' USL\n {usl}', color='darkred', fontsize=7.5, fontweight='bold', va='top', ha='right')
            if tgt is not None:
                ax.axvline(tgt, color='#1a7abf', lw=1.5, ls=':', zorder=3)
                ax.text(tgt, y_top * 0.75, f' TGT\n {tgt}', color='#1a7abf', fontsize=7, fontweight='bold', va='top', ha='left')

        ax.legend(handles=[Patch(facecolor=c_map[g], label=g) for g in base_grades if g in data.columns],
                  loc='upper right', fontsize=7)
        ax.set_xlim(fmin, fmax)
        ax.set_ylim(0, y_lim)
        ax.set_title(title, fontsize=10, fontweight='bold')
        add_chart_border(ax)

    # --- TABS ---
    tab1, tab2, tab3, tab4, tab5, tab6, tab7 = st.tabs([
        "📁 1. Raw Data", 
        "📋 2. Quality Yield", 
        "📈 3. Capability (SPC)", 
        "📉 4. I-MR Tracking",
        "✂️ 5. Tail Scrap",
        "🎯 6. Customer End-Use",
        "🏭 7. Production-Based Scrap & Material Stability"
    ])

    # ==========================================================
    # TASK 1: RAW DATA INSPECTION
    # ==========================================================
    with tab1:
        st.header("1. Raw Data Inspection")
        st.info("Grouped Thickness: 0.5mm (0.47, 0.50), 0.6mm (0.53-0.60), 0.8mm (0.63-0.80). Unmatched and empty rows removed.")
        
        col_m1, col_m2, col_m3 = st.columns(3)
        col_m1.metric("Valid Rows", f"{len(df):,}")
        col_m2.metric("Total Columns", len(df.columns))
        
        if 'Production_Date' in df.columns and not df['Production_Date'].isna().all():
            min_date = df['Production_Date'].min().strftime('%Y-%m-%d')
            max_date = df['Production_Date'].max().strftime('%Y-%m-%d')
            col_m3.metric("Date Range", f"{min_date} to {max_date}")
        else:
            col_m3.metric("Date Range", "N/A")
            
        st.markdown("### Filtered Data Preview")
        st.dataframe(df.head(50), use_container_width=True)

    # ==========================================================
    # TASK 2: YIELD SUMMARY
    
    # ==========================================================
    with tab2:
        st.header("2. Executive Quality Yield Summary")
        
        st.subheader("Detailed Yield by Thickness & Material")
        
        # BƯỚC 1: Chỉ lấy các cột chắc chắn tồn tại trong dữ liệu
        yield_summary = df.groupby(['Time_Group', 'Actual_Thickness', 'HR_Material'])[
            ['Total_Qty', 'Acceptable_Qty', 'Severe_Bad_Qty']
        ].sum().reset_index()
        
        yield_summary = yield_summary[yield_summary['Total_Qty'] > 0]
        
        if not yield_summary.empty:
            yield_summary['Yield (%)'] = (yield_summary['Acceptable_Qty'] / yield_summary['Total_Qty'] * 100).round(2)
            yield_summary['Defect_Rate (%)'] = (yield_summary['Severe_Bad_Qty'] / yield_summary['Total_Qty'] * 100).round(2)
            
            # BƯỚC 2: Tính toán Tỉ lệ Scrap gián tiếp (Scrap = Tổng - Đạt)
            yield_summary['Scrap_Rate (%)'] = ((yield_summary['Total_Qty'] - yield_summary['Acceptable_Qty']) / yield_summary['Total_Qty'] * 100).round(2)
            
            yield_summary['_sort'] = yield_summary['Time_Group'].apply(get_sort_key)
            yield_summary = yield_summary.sort_values(by=['_sort', 'Actual_Thickness']).drop(columns=['_sort'])

            # BƯỚC 3: Cập nhật hiển thị bảng dữ liệu (đã bỏ Scrap_Qty khỏi phần format)
            st.dataframe(
                yield_summary.style
                    .background_gradient(subset=['Yield (%)'], cmap='Greens')
                    .background_gradient(subset=['Defect_Rate (%)'], cmap='Reds')
                    .background_gradient(subset=['Scrap_Rate (%)'], cmap='Oranges') 
                    .format({
                        'Actual_Thickness': '{:.2f}', 'Total_Qty': '{:.0f}',
                        'Acceptable_Qty': '{:.0f}', 'Severe_Bad_Qty': '{:.0f}',
                        'Yield (%)': '{:.2f}%', 'Defect_Rate (%)': '{:.2f}%', 'Scrap_Rate (%)': '{:.2f}%'
                    }),
                use_container_width=True, hide_index=True
            )
        else:
            st.info("No yield data available to display in this view.")

        st.markdown("---")
        st.subheader("📊 Grade Distribution & Scrap by Time Period (%)")
        st.caption("Note: This summary table evaluates 100% of production data. Detailed charts below are filtered to specific thickness groups.")
        
        # 1. Tính toán phân bố Grade
        grade_dist = df_global_grades.groupby('Time_Group')[base_grades].sum()
        grade_dist['Total'] = grade_dist.sum(axis=1)
        
        grade_dist_display = pd.DataFrame()
        for g in base_grades:
            grade_dist_display[g] = (grade_dist[g] / grade_dist['Total'].replace(0, np.nan) * 100).fillna(0).round(1)
            
        # 2. TÍNH TỈ LỆ SCRAP THEO ĐÚNG LOGIC TỪ TAB 5
        COIL_ID_COL = '鋼捲號碼'
        if LEN_COL in df.columns and SCRAP_COL in df.columns and COIL_ID_COL in df.columns:
            df_temp = df.copy()
            df_temp[COIL_ID_COL] = df_temp[COIL_ID_COL].astype(str).str.strip().replace(['nan', 'None', '', 'NaN'], np.nan)
            
            missing_mask = df_temp[COIL_ID_COL].isna()
            if missing_mask.any():
                df_temp.loc[missing_mask, COIL_ID_COL] = [f"UNKNOWN_{i}" for i in df_temp[missing_mask].index]

            scrap_totals = df_temp.groupby(['Time_Group', COIL_ID_COL])[SCRAP_COL].sum().reset_index()
            
            # Xử lý drop_duplicates an toàn
            sort_cols = ['Time_Group', 'Production_Date'] if 'Production_Date' in df_temp.columns else ['Time_Group']
            first_occurrence = df_temp.sort_values(sort_cols).drop_duplicates(subset=['Time_Group', COIL_ID_COL], keep='first')
            
            df_scrap_master_temp = first_occurrence[['Time_Group', COIL_ID_COL, LEN_COL]].merge(
                scrap_totals, on=[COIL_ID_COL, 'Time_Group']
            )

            scrap_by_period_temp = df_scrap_master_temp.groupby('Time_Group').agg(
                Total_Length=(LEN_COL, 'sum'),
                Total_Scrap=(SCRAP_COL, 'sum')
            )
            
            # Tính tỉ lệ % (Làm tròn 2 chữ số thập phân để ra được 10.41%)
            scrap_by_period_temp['Scrap_Rate'] = np.where(
                scrap_by_period_temp['Total_Length'] > 0,
                (scrap_by_period_temp['Total_Scrap'] / scrap_by_period_temp['Total_Length'] * 100),
                0
            ).round(2)
            
            # Gắn vào bảng hiển thị chính
            grade_dist_display = grade_dist_display.join(scrap_by_period_temp['Scrap_Rate'])
        else:
            grade_dist_display['Scrap_Rate'] = 0.0
            
        grade_dist_display['Scrap_Rate'] = grade_dist_display['Scrap_Rate'].fillna(0)
        
        # 3. Sắp xếp dữ liệu
        grade_dist_display['_sort'] = grade_dist_display.index.map(get_sort_key)
        grade_dist_display = grade_dist_display.sort_values('_sort').drop(columns=['_sort'])
        
        # 4. Chuyển đổi định dạng thêm dấu '%'
        grade_dist_pct_str = grade_dist_display.copy()
        for col in grade_dist_display.columns:
            if col == 'Scrap_Rate':
                grade_dist_pct_str[col] = grade_dist_display[col].map(lambda x: f"{x:.2f}%") 
            else:
                grade_dist_pct_str[col] = grade_dist_display[col].map(lambda x: f"{x:.1f}%")

        # 5. Tạo HTML hiển thị bảng
        header_color = "#1a3a5c"
        alt_row_color = "#dce6f1"
        html = f"""
        <style>
        .grade-table {{ width:100%; border-collapse:collapse; font-family:sans-serif; font-size:14px; margin-bottom:24px; }}
        .grade-table th {{ background-color:{header_color}; color:white; padding:10px 16px; text-align:center; }}
        .grade-table td {{ padding:9px 16px; text-align:center; border-bottom:1px solid #ccc; }}
        .grade-table tr:nth-child(odd) td {{ background-color:{alt_row_color}; }}
        .grade-table tr:nth-child(even) td {{ background-color:#ffffff; }}
        .grade-table tr:hover td {{ background-color:#b8cce4; }}
        </style>
        <table class="grade-table">
            <thead>
                <tr>
                    <th>Time Period</th>
                    {''.join(f'<th>{g}</th>' for g in base_grades)}
                    <th style="background-color:#e67e22;">Scrap Rate</th>
                </tr>
            </thead>
            <tbody>
        """
        for period, row in grade_dist_pct_str.iterrows():
            html += "<tr>"
            html += f"<td><b>{period}</b></td>"
            
            for g in base_grades:
                val = float(row[g].replace('%', ''))
                if g in ['B+', 'B'] and val > 1.0:
                    html += f'<td style="color:#c00000;font-weight:bold">{row[g]}</td>'
                elif g in ['A-B+', 'A-B'] and val > 30.0:
                    html += f'<td style="color:#2e7d32;font-weight:bold">{row[g]}</td>'
                else:
                    html += f'<td>{row[g]}</td>'
            
            html += f'<td style="color:#d35400; font-weight:bold;">{row["Scrap_Rate"]}</td>'
            html += "</tr>"
            
        html += "</tbody></table>"
        st.markdown(html, unsafe_allow_html=True)
        # --- CODE THÊM NÚT TẢI BẢNG DỮ LIỆU ---
        # Tạo bản sao của bảng dữ liệu đang hiển thị và đặt tên cho cột thời gian
        export_df = grade_dist_pct_str.copy()
        export_df.index.name = "Time Period"

        # Chuyển đổi dữ liệu sang định dạng CSV (dùng utf-8-sig để không bị lỗi font chữ)
        csv_data = export_df.to_csv(index=True).encode('utf-8-sig')

        # Tạo nút Download
        st.download_button(
            label="📥 Tải bảng dữ liệu này (CSV)",
            data=csv_data,
            file_name="grade_distribution_and_scrap.csv",
            mime="text/csv",
            key="dl_grade_scrap"
        )
    # ==========================================================
    # TASK 3: DISTRIBUTION & PROCESS CAPABILITY (SPC)
    # ==========================================================
    with tab3:
        st.header("3. Distribution & Process Capability (SPC)")
        st.info("Visualizing mechanical property distribution. Capability indices (Cp, Cpk) apply from Q4 2025 onwards. Limit calculations strictly based on grades A or B (A-B+, A-B).")

        ordered_periods = sorted(df['Time_Group'].unique(), key=get_sort_key)
        thickness_list = sorted(df['Actual_Thickness'].dropna().unique())
        
        cap_summary_rows = []
        for _p in ordered_periods:
            if _p == "2025 (Full Year)": continue 
            _dfp = df[df['Time_Group'] == _p]
            _dfp_valid = _dfp[_dfp['Valid_Qty'] > 0]
            
            for _t in thickness_list:
                _dft = _dfp_valid[_dfp_valid['Actual_Thickness'] == _t]
                if _dft.empty: continue
                for _f in ['YS', 'TS', 'EL', 'YPE']:
                    if _f in _dft.columns:
                        row = build_capability_summary(_dft, _f, _p, _t)
                        if row: cap_summary_rows.append(row)

        if cap_summary_rows:
            cap_df = pd.DataFrame(cap_summary_rows)
            
            st.markdown("### 📋 Detailed Capability Log (Per Thickness - Grades A-B+, A-B Only)")
            def color_cpk_cell(val):
                if pd.isna(val) or val == 'N/A' or val == '': return ''
                try:
                    c = cpk_color(float(val))
                    return f'background-color:{c};color:white;font-weight:bold;text-align:center'
                except:
                    return ''

            fmt = {
                'Thickness': '{:.2f}',
                'Mean': '{:.2f}', 'Std': '{:.3f}', 
                'Cp': lambda v: f'{v:.3f}' if pd.notnull(v) else '—', 
                'Cpk': lambda v: f'{v:.3f}' if pd.notnull(v) else '—',
                'Ca (%)': lambda v: f'{v:.1f}%' if pd.notnull(v) else '—',
                'LSL': lambda v: str(v) if pd.notnull(v) else '—',
                'USL': lambda v: str(v) if pd.notnull(v) else '—',
            }
            st.dataframe(
                cap_df.style
                .map(color_cpk_cell, subset=['Cpk'])
                .format(fmt, na_rep='—'),
                use_container_width=True, hide_index=True
            )
            st.markdown("---")

        for period in ordered_periods:
            if period == "2025 (Full Year)": continue
            df_p = df[df['Time_Group'] == period]
            if df_p.empty: continue
            
            st.markdown(f"## 📅 Period: **{period}**")
            
            ov_y = get_shared_y(df_p, ['YS', 'TS', 'EL', 'YPE'])
            st.markdown(f"#### 🌐 Overall Summary (All Thicknesses)")
            cols = st.columns(2)
            for idx, f in enumerate([x for x in ['YS', 'TS', 'EL', 'YPE'] if x in df_p.columns]):
                with cols[idx % 2]:
                    df_p_valid = df_p[df_p['Valid_Qty'] > 0]
                    vals_all = df_p_valid[f].dropna().values
                    render_capability_badge(calc_capability(vals_all, f, period, 'Overall'), f, period, 'Overall')
                    
                    fig, ax = plt.subplots(figsize=(8, 4.5))
                    plot_dist(ax, df_p, f, f"{f} (Overall - {period})", ov_y, period, 'Overall')
                    fig.tight_layout()
                    st.pyplot(fig)
                    plt.close(fig)  # FIX: Ngăn sập RAM
            
            for thick in thickness_list:
                df_t = df_p[df_p['Actual_Thickness'] == thick]
                if df_t.empty: continue
                
                st.markdown(f"#### 📏 Thickness: **{thick}mm**")
                ly = get_shared_y(df_t, ['YS', 'TS', 'EL', 'YPE'])
                tcols = st.columns(2)
                
                for idx, f in enumerate([x for x in ['YS', 'TS', 'EL', 'YPE'] if x in df_t.columns]):
                    with tcols[idx % 2]:
                        df_t_valid = df_t[df_t['Valid_Qty'] > 0]
                        vals_t = df_t_valid[f].dropna().values
                        render_capability_badge(calc_capability(vals_t, f, period, thick), f, period, thick)
                        
                        fig, ax = plt.subplots(figsize=(8, 4.5))
                        plot_dist(ax, df_t, f, f"{f} (Thick:{thick} - {period})", ly, period, thick)
                        fig.tight_layout()
                        st.pyplot(fig)
                        plt.close(fig)  # FIX: Ngăn sập RAM
            st.markdown("---")

    # ==========================================================
    # ==========================================================
    # TASK 4: POST-CONTROL TRACKING (I-MR CHARTS)
    # ==========================================================
    with tab4:
        st.header("4. Post-Control Tracking (I-MR Charts)")
        st.info("Tracking process stability for production from 2026 onwards. Limits and charts calculated ONLY based on grades A or B (A-B+, A-B).")
        
        df_t4 = df[df['Production_Date'].dt.year >= 2026].copy()
        df_t4 = df_t4[df_t4['Valid_Qty'] > 0]
        
        if df_t4.empty:
            st.warning("No valid data available for 2026 onwards.")
        else:
            thicknesses = ['Overall'] + sorted(df_t4['Actual_Thickness'].dropna().unique().tolist())
            t4_thick = st.selectbox("Select Thickness Category:", thicknesses)
            
            if t4_thick != 'Overall':
                plot_df_base = df_t4[df_t4['Actual_Thickness'] == t4_thick]
            else:
                plot_df_base = df_t4
                
            for t4_feat in ['YS', 'TS', 'EL', 'YPE']:
                if t4_feat not in plot_df_base.columns:
                    continue
                    
                plot_df = plot_df_base.sort_values('Production_Date').dropna(subset=[t4_feat])
                
                if len(plot_df) < 2:
                    continue
                    
                st.markdown("---")
                st.subheader(f"Feature: {t4_feat}")
                
                vals_all = plot_df[t4_feat].values
                cap_data = calc_capability(vals_all, t4_feat, '2026-01', t4_thick)
                
                render_capability_badge(cap_data, t4_feat, '2026-01', t4_thick)
                
                dates = plot_df['Production_Date'].dt.strftime('%Y-%m-%d')
                vals = plot_df[t4_feat].values
                
                mean_v = np.mean(vals)
                mr = np.abs(np.diff(vals))
                mr_mean = np.mean(mr)
                
                ucl_i = mean_v + 2.66 * mr_mean
                lcl_i = max(0, mean_v - 2.66 * mr_mean)
                ucl_mr = 3.267 * mr_mean
                
                fig_imr, (ax_i, ax_mr) = plt.subplots(2, 1, figsize=(12, 8), gridspec_kw={'height_ratios': [2, 1]})
                
                # I-Chart
                ax_i.plot(vals, marker='o', color='#1f77b4', linestyle='-', linewidth=1.5, markersize=5)
                
                # Statistical Control Limits
                ax_i.axhline(mean_v, color='green', linestyle='--', label=f'Mean: {mean_v:.2f}')
                ax_i.axhline(ucl_i, color='red', linestyle='--', label=f'UCL: {ucl_i:.2f}')
                ax_i.axhline(lcl_i, color='red', linestyle='--', label=f'LCL: {lcl_i:.2f}')
                
                # User Input Specification Limits (USL/LSL/Target) & Mill Range
                spec = GLOBAL_SPECS.get(t4_thick, {}).get(t4_feat, {}) if t4_thick != 'Overall' else {}
                lsl = spec.get('min')
                usl = spec.get('max')
                tgt = spec.get('target')

                if lsl is not None:
                    ax_i.axhline(lsl, color='darkred', linestyle='-', linewidth=2, label=f'LSL (Spec Min): {lsl}')
                if usl is not None:
                    ax_i.axhline(usl, color='darkred', linestyle='-', linewidth=2, label=f'USL (Spec Max): {usl}')
                if tgt is not None:
                    ax_i.axhline(tgt, color='blue', linestyle=':', linewidth=1.5, label=f'Target: {tgt}')
                
                # Catch out of bounds errors
                out_condition = (vals > ucl_i) | (vals < lcl_i)
                
                if usl is not None:
                    out_condition = out_condition | (vals > usl)
                if lsl is not None:
                    out_condition = out_condition | (vals < lsl)

                out_i = np.where(out_condition)[0]
                
                if len(out_i) > 0:
                    ax_i.scatter(out_i, vals[out_i], color='red', zorder=5, s=50, label="Out of Bounds (UCL/LCL/Spec)")
                    
                all_y_vals = [v for v in [np.max(vals), np.min(vals), ucl_i, lcl_i, usl, lsl] if v is not None]
                y_max = max(all_y_vals) if all_y_vals else 10
                y_min = min(all_y_vals) if all_y_vals else 0
                y_pad = (y_max - y_min) * 0.4 if (y_max - y_min) != 0 else 1
                ax_i.set_ylim(y_min - y_pad*0.4, y_max + y_pad*1.2)
                
                ax_i.set_title(f"Individual (I) Chart - {t4_feat} ({t4_thick}mm)", fontweight='bold')
                ax_i.set_ylabel("Value")
                
                ax_i.legend(bbox_to_anchor=(1.01, 1), loc='upper left', ncol=1, fontsize=8)
                if 'add_chart_border' in globals():
                    add_chart_border(ax_i)
                ax_i.set_xticks([]) 
                
                # MR-Chart
                ax_mr.plot(range(1, len(vals)), mr, marker='o', color='#ff7f0e', linestyle='-', linewidth=1.5, markersize=5)
                ax_mr.axhline(mr_mean, color='green', linestyle='--', label=f'MR Mean: {mr_mean:.2f}')
                ax_mr.axhline(ucl_mr, color='red', linestyle='--', label=f'UCL: {ucl_mr:.2f}')
                
                out_mr = np.where(mr > ucl_mr)[0]
                if len(out_mr) > 0:
                    ax_mr.scatter(out_mr + 1, mr[out_mr], color='red', zorder=5, s=50)
                    
                ax_mr.set_title("Moving Range (MR) Chart", fontweight='bold')
                ax_mr.set_ylabel("Range")
                
                ax_mr.legend(bbox_to_anchor=(1.01, 1), loc='upper left', fontsize=8)
                if 'add_chart_border' in globals():
                    add_chart_border(ax_mr)
                
                step = max(1, len(vals) // 15)
                ax_mr.set_xticks(range(0, len(vals), step))
                ax_mr.set_xticklabels(dates.iloc[::step], rotation=45, ha='right')
                
                fig_imr.tight_layout()
                st.pyplot(fig_imr)
                
                # ==========================================
                # PPTX EXPORT FEATURE
                # ==========================================
                img_stream = io.BytesIO()
                fig_imr.savefig(img_stream, format='png', bbox_inches='tight', dpi=300)
                img_stream.seek(0)
                
                try:
                    from pptx import Presentation
                    from pptx.util import Inches
                    
                    prs = Presentation()
                    blank_slide_layout = prs.slide_layouts[6] # Layout 6 is blank
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Insert the image into the slide. Dimensions are tailored to fit standard widescreen.
                    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.5), width=Inches(9))
                    
                    pptx_stream = io.BytesIO()
                    prs.save(pptx_stream)
                    pptx_stream.seek(0)
                    
                    st.download_button(
                        label=f"📥 Download {t4_feat} Chart (.pptx)",
                        data=pptx_stream,
                        file_name=f"IMR_Chart_{t4_feat}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary",
                        use_container_width=True,
                        key=f"dl_ppt_{t4_feat}_{t4_thick}"
                    )
                except ImportError:
                    st.error("Missing dependency. Please run `pip install python-pptx` to enable PowerPoint export.")

                plt.close(fig_imr)  # Ngăn chặn tràn bộ nhớ RAM

    # ==========================================================
    # TASK 5: TAIL SCRAP & HYBRID TREND
    # ==========================================================
    import io  # Đảm bảo đã import thư viện này ở đầu file

    with tab5:
        st.header("5. Tail Scrap & Length Rejection Analysis")
        
        COIL_ID_COL = '鋼捲號碼'

        if LEN_COL in df.columns and SCRAP_COL in df.columns:
            # Drop full year summary to avoid noise
            df_t5 = df[df['Time_Group'] != "2025 (Full Year)"].copy()
            
            df_t5[COIL_ID_COL] = df_t5[COIL_ID_COL].astype(str).str.strip().replace(['nan', 'None', '', 'NaN'], np.nan)
            
            missing_mask = df_t5[COIL_ID_COL].isna()
            if missing_mask.any():
                df_t5.loc[missing_mask, COIL_ID_COL] = [f"UNKNOWN_{i}" for i in df_t5[missing_mask].index]

            scrap_totals = df_t5.groupby(['Time_Group', COIL_ID_COL])[SCRAP_COL].sum().reset_index()
            first_occurrence = df_t5.sort_values(['Time_Group', 'Production_Date']).drop_duplicates(subset=['Time_Group', COIL_ID_COL], keep='first')
            
            df_scrap_master = first_occurrence[['Time_Group', COIL_ID_COL, LEN_COL, 'Actual_Thickness', 'HR_Material', 'Production_Date']].merge(
                scrap_totals, on=[COIL_ID_COL, 'Time_Group']
            )

            # --- 1. HYBRID TREND LINE ---
            st.subheader("Rejection Rate Trend (%)")
            
            trend_data = df_scrap_master.groupby('Time_Group').agg(
                Input_Length=(LEN_COL, 'sum'),
                Total_Scrap=(SCRAP_COL, 'sum')
            ).reset_index()
            
            trend_data = trend_data[trend_data['Time_Group'] != 'Unknown']
            
            trend_data['Rejection_Rate (%)'] = np.where(
                trend_data['Input_Length'] > 0,
                (trend_data['Total_Scrap'] / trend_data['Input_Length'] * 100),
                0
            ).round(2)
            
            trend_data['_sort'] = trend_data['Time_Group'].apply(get_sort_key)
            trend_data = trend_data.sort_values('_sort').drop(columns=['_sort'])

            fig_trend, ax_trend = plt.subplots(figsize=(14, 5))
            if not trend_data.empty:
                ax_trend.plot(trend_data['Time_Group'], trend_data['Rejection_Rate (%)'], 
                              marker='o', markersize=8, markeredgecolor='white', markeredgewidth=1.5,
                              linestyle='-', color='#1f77b4', linewidth=3, label='Rejection Rate %')
                
                ax_trend.fill_between(trend_data['Time_Group'], trend_data['Rejection_Rate (%)'], 
                                      color='#1f77b4', alpha=0.1)

                y_max = trend_data['Rejection_Rate (%)'].max()
                ax_trend.set_ylim(0, y_max * 1.35 + 0.5 if not trend_data.empty else 10)
                
                ax_trend.set_title("Rejection Rate Trend", fontweight='bold', fontsize=15, pad=15, color='#333')
                ax_trend.set_ylabel("Rejection Rate (%)", fontweight='bold', color='#555')
                ax_trend.set_xlabel("")
                
                ax_trend.grid(axis='y', linestyle='--', alpha=0.5)
                ax_trend.grid(axis='x', visible=False)
                
                for i, val in enumerate(trend_data['Rejection_Rate (%)']):
                    ax_trend.annotate(f'{val:.2f}%', 
                                      xy=(i, val), 
                                      xytext=(0, 8), 
                                      textcoords="offset points", 
                                      ha='center', va='bottom', 
                                      fontsize=10, fontweight='bold', color='#222',
                                      bbox=dict(boxstyle="round,pad=0.3", fc="white", ec="none", alpha=0.8))
                
                add_chart_border(ax_trend)
                plt.xticks(rotation=40, ha='right', fontsize=10)
                fig_trend.tight_layout()
                
            st.pyplot(fig_trend)
            
            # --- THÊM NÚT TẢI ẢNH 1 ---
            buf_trend = io.BytesIO()
            fig_trend.savefig(buf_trend, format="png", bbox_inches="tight", dpi=300)
            buf_trend.seek(0)
            st.download_button(
                label="📥 Tải biểu đồ Trend về máy (PNG)",
                data=buf_trend,
                file_name="rejection_rate_trend.png",
                mime="image/png",
                key="dl_trend"
            )
            plt.close(fig_trend)  # FIX: Ngăn sập RAM

            # --- 2. PERIOD SUMMARY & CHART ---
            st.markdown("---")
            st.subheader("Scrap Rate by Time Period")
            scrap_by_period = df_scrap_master.groupby('Time_Group').agg(
                Total_Length=(LEN_COL, 'sum'),
                Total_Scrap=(SCRAP_COL, 'sum'),
                Coil_Count=(COIL_ID_COL, 'count')
            ).reset_index()
            
            scrap_by_period['Scrap_Rate (%)'] = np.where(
                scrap_by_period['Total_Length'] > 0,
                (scrap_by_period['Total_Scrap'] / scrap_by_period['Total_Length'] * 100),
                0
            ).round(2)
            
            scrap_by_period['_sort'] = scrap_by_period['Time_Group'].apply(get_sort_key)
            scrap_by_period = scrap_by_period.sort_values('_sort').drop(columns=['_sort'])
            
            fig_p, ax_p = plt.subplots(figsize=(10, 4))
            if not scrap_by_period.empty:
                ax_p.bar(scrap_by_period['Time_Group'], scrap_by_period['Scrap_Rate (%)'], color='#e74c3c', edgecolor='white')
                ax_p.set_title("Tail Scrap Rate (%) by Time Period", fontweight='bold')
                ax_p.set_ylabel("Scrap Rate (%)")
                ax_p.set_xlabel("")
                ax_p.set_ylim(0, scrap_by_period['Scrap_Rate (%)'].max() * 1.2 + 0.1)
                for i, val in enumerate(scrap_by_period['Scrap_Rate (%)']):
                    ax_p.annotate(f"{val:.2f}%", xy=(i, val), xytext=(0, 5), textcoords="offset points", ha='center', va='bottom', fontweight='bold')
                
                add_chart_border(ax_p)
                plt.xticks(rotation=30, ha='right')
                fig_p.tight_layout()
                
            st.pyplot(fig_p)
            
            # --- THÊM NÚT TẢI ẢNH 2 ---
            buf_p = io.BytesIO()
            fig_p.savefig(buf_p, format="png", bbox_inches="tight", dpi=300)
            buf_p.seek(0)
            st.download_button(
                label="📥 Tải biểu đồ Scrap Rate (PNG)",
                data=buf_p,
                file_name="scrap_rate_by_period.png",
                mime="image/png",
                key="dl_period"
            )
            plt.close(fig_p)  # FIX: Ngăn sập RAM

            st.dataframe(
                scrap_by_period.style.background_gradient(subset=['Scrap_Rate (%)'], cmap='Reds')
                .format({'Total_Length': '{:,.2f}', 'Total_Scrap': '{:,.2f}', 'Scrap_Rate (%)': '{:.2f}%'}),
                use_container_width=True, hide_index=True
            )

            # --- 3. LEVEL-BY-LEVEL DRILL DOWN & CHARTS ---
            st.markdown("---")
            st.subheader("Deep Analysis: Scrap Rate by Period / Thickness / Material")
            
            scrap_detail = df_scrap_master.groupby(['Time_Group', 'Actual_Thickness', 'HR_Material']).agg(
                Total_Length=(LEN_COL, 'sum'),
                Total_Scrap=(SCRAP_COL, 'sum'),
                Coil_Count=(COIL_ID_COL, 'count')
            ).reset_index()
            
            scrap_detail = scrap_detail[scrap_detail['Total_Length'] > 0]
            scrap_detail['Scrap_Rate (%)'] = (scrap_detail['Total_Scrap'] / scrap_detail['Total_Length'] * 100).round(2)
            
            col_t, col_m = st.columns(2)
            
            with col_t:
                st.markdown("**Scrap Rate by Period & Thickness**")
                fig_t, ax_t = plt.subplots(figsize=(8, 4))
                if not scrap_detail.empty:
                    thick_agg = scrap_detail.groupby(['Time_Group', 'Actual_Thickness']).agg(
                        T_Len=('Total_Length', 'sum'),
                        T_Scrap=('Total_Scrap', 'sum')
                    )
                    thick_agg['Rate'] = np.where(thick_agg['T_Len'] > 0, (thick_agg['T_Scrap'] / thick_agg['T_Len']) * 100, 0)
                    pivot_t = thick_agg['Rate'].unstack()
                    
                    if not pivot_t.empty:
                        pivot_t.plot(kind='bar', ax=ax_t, colormap='tab10', edgecolor='white')
                        ax_t.legend(title="Thickness", bbox_to_anchor=(1.02, 1), loc='upper left')
                        for c in ax_t.containers:
                            labels = [f"{v.get_height():.1f}%" if v.get_height() > 0 else "" for v in c]
                            ax_t.bar_label(c, labels=labels, padding=3, fontsize=7, fontweight='bold', rotation=90)
                    
                    y_max = pivot_t.max().max() if not pivot_t.isna().all().all() else 10
                    ax_t.set_ylim(0, y_max * 1.4 + 2)
                ax_t.set_ylabel("Scrap Rate (%)")
                ax_t.set_xlabel("")
                add_chart_border(ax_t)
                plt.xticks(rotation=30, ha='right')
                fig_t.tight_layout()
                
                st.pyplot(fig_t)
                
                # --- THÊM NÚT TẢI ẢNH 3 ---
                buf_t = io.BytesIO()
                fig_t.savefig(buf_t, format="png", bbox_inches="tight", dpi=300)
                buf_t.seek(0)
                st.download_button(
                    label="📥 Tải ảnh Thickness",
                    data=buf_t,
                    file_name="scrap_thickness.png",
                    mime="image/png",
                    key="dl_thick"
                )
                plt.close(fig_t)  # FIX: Ngăn sập RAM

            with col_m:
                st.markdown("**Scrap Rate by Period & Material**")
                fig_m, ax_m = plt.subplots(figsize=(8, 4))
                if not scrap_detail.empty:
                    mat_agg = scrap_detail.groupby(['Time_Group', 'HR_Material']).agg(
                        M_Len=('Total_Length', 'sum'),
                        M_Scrap=('Total_Scrap', 'sum')
                    )
                    mat_agg['Rate'] = np.where(mat_agg['M_Len'] > 0, (mat_agg['M_Scrap'] / mat_agg['M_Len']) * 100, 0)
                    pivot_m = mat_agg['Rate'].unstack()
                    
                    if not pivot_m.empty:
                        pivot_m.plot(kind='bar', ax=ax_m, colormap='tab10', edgecolor='white')
                        ax_m.legend(title="Material", bbox_to_anchor=(1.02, 1), loc='upper left')
                        for c in ax_m.containers:
                            labels = [f"{v.get_height():.1f}%" if v.get_height() > 0 else "" for v in c]
                            ax_m.bar_label(c, labels=labels, padding=3, fontsize=7, fontweight='bold', rotation=90)
                    
                    y_max = pivot_m.max().max() if not pivot_m.isna().all().all() else 10
                    ax_m.set_ylim(0, y_max * 1.4 + 2)
                ax_m.set_ylabel("Scrap Rate (%)")
                ax_m.set_xlabel("")
                add_chart_border(ax_m)
                plt.xticks(rotation=30, ha='right')
                fig_m.tight_layout()
                
                st.pyplot(fig_m)
                
                # --- THÊM NÚT TẢI ẢNH 4 ---
                buf_m = io.BytesIO()
                fig_m.savefig(buf_m, format="png", bbox_inches="tight", dpi=300)
                buf_m.seek(0)
                st.download_button(
                    label="📥 Tải ảnh Material",
                    data=buf_m,
                    file_name="scrap_material.png",
                    mime="image/png",
                    key="dl_mat"
                )
                plt.close(fig_m)  # FIX: Ngăn sập RAM

            scrap_detail['_sort'] = scrap_detail['Time_Group'].apply(get_sort_key)
            scrap_detail = scrap_detail.sort_values(by=['_sort', 'Actual_Thickness']).drop(columns=['_sort'])

            st.dataframe(
                scrap_detail.style.background_gradient(subset=['Scrap_Rate (%)'], cmap='Oranges')
                .format({'Actual_Thickness': '{:.2f}', 'Total_Length': '{:,.2f}', 'Total_Scrap': '{:,.2f}', 'Scrap_Rate (%)': '{:.2f}%'}),
                use_container_width=True, hide_index=True
            )

        else:
            st.warning("Required columns ('實測長度' or '尾料剔退') not found in the file.")
            
    # ==========================================================
    # ==========================================================
    # TASK 6: CUSTOMER END-USE ANALYSIS & MACHINE TRANSITION
    # ==========================================================
    with tab6:
        st.header("6. Customer End-Use Analysis & Machine Transition")
        st.info("Customer End-Use Root Cause Verification System: Evaluating material stability vs. machine impact.")
        
        possible_usage_cols = ['使用日期', '使用月份', 'Usage Date', 'Usage Month']
        USAGE_COL = next((c for c in possible_usage_cols if c in df.columns), None) 
        COIL_ID_COL = '鋼捲號碼'
        
        # Safely identify the weight column
        possible_wt_cols = ['重量', 'Weight', 'WT', 'Net_Weight', 'Net Weight']
        WT_COL = next((c for c in possible_wt_cols if c in df.columns), 'Weight')

        if USAGE_COL and COIL_ID_COL in df.columns and LEN_COL in df.columns and SCRAP_COL in df.columns: 
            df_t6 = df[df[LEN_COL] > 0].copy() 
            df_t6[COIL_ID_COL] = df_t6[COIL_ID_COL].astype(str).str.strip()
            df_t6 = df_t6[df_t6[COIL_ID_COL] != 'nan']

            if WT_COL in df_t6.columns:
                df_t6[WT_COL] = pd.to_numeric(df_t6[WT_COL], errors='coerce').fillna(0)

            # Vectorized Date Parsing
            if not pd.api.types.is_datetime64_any_dtype(df_t6[USAGE_COL]):
                df_t6['Usage_Date'] = pd.to_datetime(df_t6[USAGE_COL].astype(str).str.strip(), dayfirst=True, errors='coerce')
            else:
                df_t6['Usage_Date'] = df_t6[USAGE_COL]

            df_t6 = df_t6.dropna(subset=['Usage_Date'])

            # =========================================================================
            # CRITICAL LOGIC FIX: COIL-LEVEL HYBRID ASSIGNMENT
            # 1. Completion Month (Usage_Month): Get the LAST record (keep='last')
            # 2. Length & Weight: Get from the FIRST record (keep='first')
            # 3. Scrap Amount: Aggregate sum of all cuts
            # =========================================================================
            df_sorted = df_t6.sort_values('Usage_Date')
            
            # Base DataFrame: Assign coil to the final completion month
            df_coil = df_sorted.drop_duplicates(subset=[COIL_ID_COL], keep='last').copy()
            
            # Map Length & Weight from the very first appearance to avoid double counting
            df_first = df_sorted.drop_duplicates(subset=[COIL_ID_COL], keep='first')
            df_coil[LEN_COL] = df_coil[COIL_ID_COL].map(df_first.set_index(COIL_ID_COL)[LEN_COL].to_dict())
            
            if WT_COL in df_t6.columns:
                df_coil[WT_COL] = df_coil[COIL_ID_COL].map(df_first.set_index(COIL_ID_COL)[WT_COL].to_dict())
                
            # Aggregate total scrap across all records for that specific coil
            df_coil[SCRAP_COL] = df_coil[COIL_ID_COL].map(df_sorted.groupby(COIL_ID_COL)[SCRAP_COL].sum().to_dict())
            # =========================================================================

            # Usage group formatting logic
            def format_usage_group(d):
                if d.year <= 2024:
                    return "2024 (Full Year)"
                elif d.year == 2025:
                    if d <= pd.Timestamp(2025, 6, 28):
                        return "2025 H1 (Until 06/28)"
                    elif pd.Timestamp(2025, 6, 29) <= d <= pd.Timestamp(2025, 9, 30):
                        return "2025 Q3 (06/29 - 09/30)"
                    else:
                        return d.strftime('%Y-%m') # 2025 Q4 and beyond
                else:
                    return d.strftime('%Y-%m') # 2026 and beyond
            
            df_coil['Usage_Month'] = df_coil['Usage_Date'].apply(format_usage_group)

            if df_coil.empty:
                st.warning("No usage data available.")
            else:
                cutoff_date = pd.to_datetime('2026-04-01')
                df_coil['Machine_Status'] = np.where(df_coil['Usage_Date'] >= cutoff_date, 'New Machine (>= Apr 2026)', 'Old Machine (< Apr 2026)')

                props_cols = [c for c in ['YS', 'TS', 'EL', 'YPE'] if c in df_coil.columns]
                if props_cols:
                    df_coil[props_cols] = df_coil[props_cols].apply(pd.to_numeric, errors='coerce')
                    df_coil[props_cols] = df_coil[props_cols].where(df_coil[props_cols] > 0, np.nan)

                # ==========================================
                # Monthly Scrap & Material Stability Analysis
                # ==========================================
                st.subheader("Monthly Scrap & Material Stability Analysis")
                st.caption("Verifying if the spike in scrap correlates with material instability.")

                macro_df = df_coil.groupby('Usage_Month').agg(
                    Total_Length=(LEN_COL, 'sum'), 
                    Total_Scrap=(SCRAP_COL, 'sum'),
                    Avg_YS=('YS', 'mean') if 'YS' in props_cols else (LEN_COL, 'count'),
                    Avg_TS=('TS', 'mean') if 'TS' in props_cols else (LEN_COL, 'count'),
                    Avg_EL=('EL', 'mean') if 'EL' in props_cols else (LEN_COL, 'count'),
                    Avg_YPE=('YPE', 'mean') if 'YPE' in props_cols else (LEN_COL, 'count')
                ).reset_index().sort_values('Usage_Month')
                
                macro_df['Scrap_Rate (%)'] = np.where(macro_df['Total_Length'] > 0, (macro_df['Total_Scrap'] / macro_df['Total_Length']) * 100, 0).round(2)

                row1_cols = st.columns(2)
                row2_cols = st.columns(2)
                cols = row1_cols + row2_cols 
                
                features = [('Avg_YS', 'Actual YS', '#1f77b4'), 
                            ('Avg_TS', 'Actual TS', '#2ca02c'), 
                            ('Avg_EL', 'Actual EL', '#9467bd'),
                            ('Avg_YPE', 'Actual YPE', '#ff7f0e')]

                for idx, (col_name, label, color) in enumerate(features):
                    if col_name not in macro_df.columns or (macro_df[col_name] == macro_df['Total_Length']).all():
                        continue 

                    with cols[idx]:
                        fig_exec, ax1 = plt.subplots(figsize=(6, 4))
                        color1 = '#d62728' 
                        ax1.set_ylabel('Scrap Rate (%)', color=color1, fontweight='bold', fontsize=9)
                        ax1.plot(macro_df['Usage_Month'], macro_df['Scrap_Rate (%)'], color=color1, marker='o', linewidth=2, label='Scrap Rate')
                        ax1.tick_params(axis='y', labelcolor=color1, labelsize=8)
                        
                        max_scrap = macro_df['Scrap_Rate (%)'].max()
                        ax1.set_ylim(-0.5, max_scrap * 1.35 + 1 if max_scrap > 0 else 5)

                        ax2 = ax1.twinx()
                        if col_name in macro_df.columns:
                            feat_valid = macro_df[col_name].dropna()
                            if not feat_valid.empty:
                                ax2.set_ylabel(label, color=color, fontweight='bold', fontsize=9)
                                ax2.plot(macro_df['Usage_Month'], macro_df[col_name], color=color, marker='s', linestyle='--', linewidth=2, alpha=0.8, label=label)
                                ax2.tick_params(axis='y', labelcolor=color, labelsize=8)
                                feat_min, feat_max = feat_valid.min(), feat_valid.max()
                                padding = (feat_max - feat_min) * 0.15 if feat_max > feat_min else feat_min * 0.1
                                ax2.set_ylim(feat_min - padding, feat_max + padding)

                        # Adjust reference line formatting if necessary
                        plt.title(f"Scrap vs {label}", fontweight='bold', fontsize=10)
                        ax1.set_xticklabels(macro_df['Usage_Month'], rotation=45, ha='right', fontsize=8)
                        
                        if 'add_chart_border' in globals():
                            add_chart_border(ax1) 
                            
                        fig_exec.tight_layout()
                        st.pyplot(fig_exec)
                        
                        buf = io.BytesIO()
                        fig_exec.savefig(buf, format="png", dpi=300, bbox_inches="tight")
                        buf.seek(0)
                        st.download_button(
                            label=f"📸 Download {label} Chart",
                            data=buf,
                            file_name=f"Scrap_vs_{col_name}.png",
                            mime="image/png",
                            type="primary",
                            use_container_width=True,
                            key=f"dl_chart_{idx}" 
                        )
                        plt.close(fig_exec) 

                st.markdown("<div style='text-align: center; color: #c00000; font-weight: bold; font-size: 14px; margin-bottom: 20px;'>Logic: If Scrap increases but YS/TS/EL/YPE is stable ➡️ Issue is with the Customer's Machine.</div>", unsafe_allow_html=True)
                st.markdown("---")
                
                # ==========================================
                # Production vs Usage Quality Matrix
                # ==========================================
                st.subheader("Production vs Usage Quality Matrix (Main Chart)")
                st.info("Evaluates Material Stability, Inventory Traceability, Machine Impact, and Quality Transition.")

                # Calculate Matrix data purely from the deduplicated df_coil
                agg_dict = {
                    'Total_Length': (LEN_COL, 'sum'), 
                    'Total_Scrap': (SCRAP_COL, 'sum'), 
                    'Total_Coils': (COIL_ID_COL, 'count') 
                }
                matrix_data = df_coil.groupby(['Usage_Month', 'Time_Group']).agg(**agg_dict).reset_index()
                
                available_grades = [g for g in base_grades if g in df_coil.columns] if 'base_grades' in globals() else []
                if available_grades:
                    grade_data = df_coil.groupby(['Usage_Month', 'Time_Group'])[available_grades].sum().reset_index()
                    matrix_data = pd.merge(matrix_data, grade_data, on=['Usage_Month', 'Time_Group'], how='left')

                matrix_data['Scrap_Rate'] = np.where(matrix_data['Total_Length'] > 0, (matrix_data['Total_Scrap'] / matrix_data['Total_Length']) * 100, 0).round(2)
                
                # ---------------------------------------------------------
                # Dynamic Custom Sort Logic
                # Groups: 1(Q/H) -> 2(Months) -> 3(Full Year)
                # ---------------------------------------------------------
                def custom_time_sort(period_str):
                    p = str(period_str)
                    year = p[:4]
                    
                    if "Full Year" in p:
                        group = "3_FullYear"
                    elif any(q in p for q in ["H1", "H2", "Q1", "Q2", "Q3", "Q4"]):
                        group = f"1_{p}"
                    elif len(p) >= 7 and "-" in p[4:8]:
                        group = f"2_{p}"
                    else:
                        group = f"4_{p}"
                        
                    return f"{year}_{group}"

                prod_periods = sorted(matrix_data['Time_Group'].unique(), key=custom_time_sort)
                usage_months = sorted(matrix_data['Usage_Month'].unique())

                # Coil Level Summary Logic (Rows/Columns boundary totals)
                prod_summary = df_coil.groupby('Time_Group').agg({
                    LEN_COL: 'sum',
                    WT_COL: 'sum' if WT_COL in df_coil.columns else lambda x: 0
                }).to_dict('index')
                
                usage_summary = df_coil.groupby('Usage_Month').agg({
                    LEN_COL: 'sum',
                    WT_COL: 'sum' if WT_COL in df_coil.columns else lambda x: 0
                }).to_dict('index')

                total_matrix_L = df_coil[LEN_COL].sum()
                total_matrix_W = df_coil[WT_COL].sum() if WT_COL in df_coil.columns else 0

                def get_color(rate):
                    if pd.isna(rate): return "#ffffff" 
                    if rate < 2.0: return "#e8f5e9" 
                    if rate < 5.0: return "#fff3e0" 
                    if rate < 10.0: return "#ffcdd2" 
                    return "#e57373" 

                matrix_dict = matrix_data.set_index(['Time_Group', 'Usage_Month']).to_dict('index')

                html_parts = [
                    "<style>",
                    ".q-matrix { width: 100%; border-collapse: collapse; font-family: sans-serif; font-size: 12px; }",
                    ".q-matrix th { background-color: #1a3a5c; color: white; padding: 10px; text-align: center; border: 1px solid #ddd; }",
                    ".q-matrix td { border: 1px solid #ccc; padding: 8px; vertical-align: top; }",
                    ".cell-title { font-size: 14px; font-weight: bold; margin-bottom: 5px; color: #111; text-align: center; border-bottom: 1px solid rgba(0,0,0,0.1); padding-bottom: 3px;}",
                    ".grade-list { list-style-type: none; padding: 0; margin: 0; line-height: 1.4; }",
                    ".grade-list li { display: flex; justify-content: space-between; }",
                    ".grade-name { font-weight: bold; color: #444; }",
                    ".summary-cell { background-color: #e3f2fd; color: #0d47a1; text-align: center; vertical-align: middle; font-weight: bold; font-size: 11px; }",
                    ".summary-header { background-color: #1565c0 !important; }",
                    "</style>",
                    "<table class='q-matrix'><thead><tr><th>Production \\ Usage</th>"
                ]
                html_parts.extend([f"<th>{m}</th>" for m in usage_months])
                html_parts.append("<th class='summary-header'>Total Output<br>(生產總量)</th>")
                html_parts.append("</tr></thead><tbody>")

                for prod in prod_periods:
                    html_parts.append(f"<tr><th style='background-color: #f1f3f5; color: #333;'>{prod}</th>")
                    for usage in usage_months:
                        row = matrix_dict.get((prod, usage))
                        if not row:
                            html_parts.append("<td style='background-color: #fafafa; color: #aaa; text-align:center; vertical-align:middle;'>No Data</td>")
                        else:
                            scrap_rate = row['Scrap_Rate']
                            bg_color = get_color(scrap_rate)
                            grade_html = []
                            total_coils = row.get('Total_Coils', 0)
                            
                            # Calculate the strict sum of grades in this exact cell to serve as 100% denominator
                            cell_total_grade = sum(row.get(g, 0) for g in available_grades) if available_grades else 0
                            
                            cell_title_html = f"<div class='cell-title'>Scrap: {scrap_rate:.1f}%<br><span style='font-size: 11px; color: #555;'>Coils: {int(total_coils)}</span></div>"

                            if cell_total_grade > 0 and available_grades:
                                for g in available_grades:
                                    g_pct = (row.get(g, 0) / cell_total_grade * 100)
                                    if g_pct > 0:
                                        color = "green" if "A" in g else "red"
                                        grade_html.append(f"<li><span class='grade-name'>{g}:</span> <span style='color:{color}'>{g_pct:.0f}%</span></li>")
                            
                            html_parts.append(f"<td style='background-color: {bg_color};'>{cell_title_html}<ul class='grade-list'>{''.join(grade_html)}</ul></td>")
                    
                    p_len = prod_summary.get(prod, {}).get(LEN_COL, 0)
                    p_wt = prod_summary.get(prod, {}).get(WT_COL, 0)
                    html_parts.append(f"<td class='summary-cell'>L: {p_len:,.0f} m<br>W: {p_wt:,.0f} kg</td>")
                    html_parts.append("</tr>")

                html_parts.append("<tr><th class='summary-header'>Total Usage<br>(客戶使用量)</th>")
                for usage in usage_months:
                    u_len = usage_summary.get(usage, {}).get(LEN_COL, 0)
                    u_wt = usage_summary.get(usage, {}).get(WT_COL, 0)
                    html_parts.append(f"<td class='summary-cell'>L: {u_len:,.0f} m<br>W: {u_wt:,.0f} kg</td>")
                
                html_parts.append(f"<td class='summary-cell' style='background-color: #bbdefb; color: #b71c1c;'>Total L: {total_matrix_L:,.0f} m<br>Total W: {total_matrix_W:,.0f} kg</td>")
                html_parts.append("</tr>")
                
                html_parts.append("</tbody></table>")

                matrix_html_str = "".join(html_parts)
                capture_component = f"""
                <!DOCTYPE html>
                <html>
                <head>
                    <script src="https://cdnjs.cloudflare.com/ajax/libs/html2canvas/1.4.1/html2canvas.min.js"></script>
                    <style>
                        body {{ font-family: sans-serif; margin: 0; padding: 0; }}
                        .btn-capture {{
                            background-color: #FF4B4B; color: white; border: none; padding: 10px;
                            border-radius: 5px; cursor: pointer; font-weight: bold; font-size: 13px;
                            margin-bottom: 10px; transition: 0.3s; width: 100%;
                        }}
                        .btn-capture:hover {{ background-color: #ff3333; }}
                    </style>
                </head>
                <body>
                    <button class="btn-capture" onclick="takeSnapshot()">📸 Download High-Resolution Matrix Chart</button>
                    <div id="matrix-container" style="background: white; padding: 10px; display: inline-block; width: 100%;">
                        {matrix_html_str}
                    </div>
                    <script>
                        function takeSnapshot() {{
                            const target = document.getElementById('matrix-container');
                            html2canvas(target, {{ scale: 3, backgroundColor: '#ffffff' }}).then(canvas => {{
                                let link = document.createElement('a');
                                link.download = 'Quality_Matrix.png';
                                link.href = canvas.toDataURL('image/png');
                                link.click();
                            }});
                        }}
                    </script>
                </body>
                </html>
                """
                components.html(capture_component, height=max(250, len(prod_periods)*65 + 100), scrolling=True)
                
                # ==========================================
                # Download Matrix to Word (Native Word Table)
                # ==========================================
                st.markdown("### 📄 Download Production Matrix Report")
                
                def create_word_report():
                    from docx import Document
                    from docx.shared import Inches, Pt, RGBColor
                    from docx.oxml.ns import nsdecls
                    from docx.oxml import parse_xml
                    from docx.enum.text import WD_ALIGN_PARAGRAPH
                    from docx.enum.table import WD_ALIGN_VERTICAL

                    doc = Document()
                    
                    section = doc.sections[0]
                    new_width, new_height = section.page_height, section.page_width
                    section.page_width = new_width
                    section.page_height = new_height
                    section.left_margin = Inches(0.5)
                    section.right_margin = Inches(0.5)

                    doc.add_heading('Production vs Usage Quality Matrix', level=1)
                    
                    def set_cell_background(cell, hex_color):
                        hex_color = hex_color.replace("#", "")
                        shading_elm = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'), hex_color))
                        cell._tc.get_or_add_tcPr().append(shading_elm)

                    cols_count = len(usage_months) + 2
                    table = doc.add_table(rows=1, cols=cols_count)
                    table.style = 'Table Grid'
                    
                    hdr_cells = table.rows[0].cells
                    hdr_cells[0].text = "Production \\ Usage"
                    set_cell_background(hdr_cells[0], "1a3a5c")
                    hdr_cells[0].paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                    hdr_cells[0].paragraphs[0].runs[0].font.bold = True
                    
                    for i, m in enumerate(usage_months):
                        hdr_cells[i+1].text = m
                        set_cell_background(hdr_cells[i+1], "1a3a5c")
                        hdr_cells[i+1].paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                        hdr_cells[i+1].paragraphs[0].runs[0].font.bold = True
                        
                    hdr_cells[-1].text = "Total Output\n(生產總量)"
                    set_cell_background(hdr_cells[-1], "1565c0")
                    hdr_cells[-1].paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                    hdr_cells[-1].paragraphs[0].runs[0].font.bold = True

                    for cell in hdr_cells:
                        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                        cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER

                    for prod in prod_periods:
                        row_cells = table.add_row().cells
                        
                        row_cells[0].text = prod
                        set_cell_background(row_cells[0], "f1f3f5")
                        row_cells[0].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                        row_cells[0].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                        row_cells[0].paragraphs[0].runs[0].font.bold = True
                        
                        for i, usage in enumerate(usage_months):
                            row = matrix_dict.get((prod, usage))
                            cell = row_cells[i+1]
                            cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                            p = cell.paragraphs[0]
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            
                            if not row:
                                run = p.add_run("No Data")
                                run.font.color.rgb = RGBColor(170, 170, 170)
                                run.font.size = Pt(8)
                                set_cell_background(cell, "fafafa")
                            else:
                                scrap_rate = row['Scrap_Rate']
                                total_coils = row.get('Total_Coils', 0)
                                bg_color = get_color(scrap_rate)
                                set_cell_background(cell, bg_color)
                                
                                # Summing up grades for denominator logic in Word Export
                                cell_total_grade = sum(row.get(g, 0) for g in available_grades) if available_grades else 0
                                
                                run_scrap = p.add_run(f"Scrap: {scrap_rate:.1f}%\nCoils: {int(total_coils)}\n")
                                run_scrap.font.bold = True
                                run_scrap.font.size = Pt(9)
                                
                                if cell_total_grade > 0 and available_grades:
                                    for g in available_grades:
                                        g_pct = (row.get(g, 0) / cell_total_grade * 100)
                                        if g_pct > 0:
                                            run_g = p.add_run(f"{g}: ")
                                            run_g.font.size = Pt(8)
                                            
                                            run_pct = p.add_run(f"{g_pct:.0f}%\n")
                                            run_pct.font.size = Pt(8)
                                            run_pct.font.bold = True
                                            if "A" in g:
                                                run_pct.font.color.rgb = RGBColor(0, 128, 0)
                                            else:
                                                run_pct.font.color.rgb = RGBColor(220, 20, 60)
                                                
                        p_len = prod_summary.get(prod, {}).get(LEN_COL, 0)
                        p_wt = prod_summary.get(prod, {}).get(WT_COL, 0)
                        cell_out = row_cells[-1]
                        set_cell_background(cell_out, "e3f2fd")
                        cell_out.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                        p_out = cell_out.paragraphs[0]
                        p_out.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run_out = p_out.add_run(f"L: {p_len:,.0f} m\nW: {p_wt:,.0f} kg")
                        run_out.font.size = Pt(8)
                        run_out.font.color.rgb = RGBColor(13, 71, 161)
                        run_out.font.bold = True

                    row_cells = table.add_row().cells
                    row_cells[0].text = "Total Usage\n(客戶使用量)"
                    set_cell_background(row_cells[0], "1565c0")
                    row_cells[0].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                    row_cells[0].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                    row_cells[0].paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                    row_cells[0].paragraphs[0].runs[0].font.bold = True
                    
                    for i, usage in enumerate(usage_months):
                        u_len = usage_summary.get(usage, {}).get(LEN_COL, 0)
                        u_wt = usage_summary.get(usage, {}).get(WT_COL, 0)
                        cell = row_cells[i+1]
                        set_cell_background(cell, "e3f2fd")
                        cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                        p = cell.paragraphs[0]
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run = p.add_run(f"L: {u_len:,.0f} m\nW: {u_wt:,.0f} kg")
                        run.font.size = Pt(8)
                        run.font.color.rgb = RGBColor(13, 71, 161)
                        run.font.bold = True
                        
                    cell_grand = row_cells[-1]
                    set_cell_background(cell_grand, "bbdefb")
                    cell_grand.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                    p_grand = cell_grand.paragraphs[0]
                    p_grand.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run_grand = p_grand.add_run(f"Total L: {total_matrix_L:,.0f} m\nTotal W: {total_matrix_W:,.0f} kg")
                    run_grand.font.size = Pt(9)
                    run_grand.font.color.rgb = RGBColor(183, 28, 28)
                    run_grand.font.bold = True

                    word_buffer = io.BytesIO()
                    doc.save(word_buffer)
                    word_buffer.seek(0)
                    return word_buffer
                    
                word_buffer = create_word_report()
                
                st.download_button(
                    label="📥 Download Native Matrix Report (.docx)",
                    data=word_buffer,
                    file_name="Production_Matrix_Report_HQ.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    type="primary",
                    use_container_width=True,
                    key="dl_matrix_word"
                )
                st.caption("Matrix Logic: Columns = Usage Month | Rows = Production Period | Background Color = Scrap Severity | Text = Quality Grade Distribution (%)")
                st.markdown("---")
                
                # ==========================================
                # Heatmap & Grade Distribution Analysis
                # ==========================================
                col_h1, col_h2 = st.columns(2)

                with col_h1:
                    st.subheader("7. Scrap Heatmap")
                    pivot_scrap = matrix_data.pivot(index='Usage_Month', columns='Time_Group', values='Scrap_Rate')
                    ordered_cols = [c for c in prod_periods if c in pivot_scrap.columns]
                    pivot_scrap = pivot_scrap[ordered_cols]
                    
                    fig_h1, ax_h1 = plt.subplots(figsize=(8, max(4, len(pivot_scrap) * 0.6)))
                    import seaborn as sns
                    sns.heatmap(pivot_scrap, annot=True, fmt=".1f", cmap="Reds", linewidths=1, linecolor='white', ax=ax_h1, annot_kws={"size": 10, "weight": "bold"})
                    ax_h1.set_ylabel("Usage Month", fontweight='bold')
                    ax_h1.set_xlabel("Production Period", fontweight='bold')
                    plt.xticks(rotation=45, ha='right')
                    fig_h1.tight_layout()
                    st.pyplot(fig_h1)
                    
                    buf_h1 = io.BytesIO()
                    fig_h1.savefig(buf_h1, format="png", dpi=300, bbox_inches="tight")
                    buf_h1.seek(0)
                    st.download_button(
                        label="📸 Download Heatmap",
                        data=buf_h1,
                        file_name="Scrap_Heatmap.png",
                        mime="image/png",
                        type="primary",
                        use_container_width=True,
                        key="dl_heatmap"
                    )
                    plt.close(fig_h1) 
                
                with col_h2:
                    st.subheader("8. Grade Distribution Analysis")
                    if available_grades:
                        # Dùng df_coil đã được làm sạch để vẽ bar chart chuẩn xác
                        grade_agg_usage = df_coil.groupby('Usage_Month')[available_grades].sum()
                        grade_pct_usage = grade_agg_usage.div(grade_agg_usage.sum(axis=1), axis=0) * 100
                        grade_pct_usage = grade_pct_usage.fillna(0)
                        
                        fig_g2, ax_g2 = plt.subplots(figsize=(8, max(4, len(pivot_scrap) * 0.6))) 
                        color_map = {'A-B+': '#2e7d32', 'A-B': '#1f77b4', 'A-B-': '#ffa726', 'B+': '#ef5350', 'B': '#c62828'}
                        plot_colors = [color_map.get(g, '#888') for g in available_grades]
                        
                        grade_pct_usage.plot(kind='bar', stacked=True, ax=ax_g2, color=plot_colors, width=0.8, edgecolor='white')
                        ax_g2.set_ylabel("Percentage (%)", fontweight='bold')
                        ax_g2.set_xlabel("Usage Month", fontweight='bold')
                        ax_g2.legend(title="Quality Grade", bbox_to_anchor=(1.02, 1), loc='upper left')
                        ax_g2.set_ylim(0, 105)
                        
                        for c in ax_g2.containers:
                            labels = [f"{v.get_height():.0f}%" if v.get_height() > 5 else "" for v in c]
                            ax_g2.bar_label(c, labels=labels, label_type='center', color='white', fontweight='bold', fontsize=9)
                            
                        plt.xticks(rotation=45, ha='right')
                        if 'add_chart_border' in globals():
                            add_chart_border(ax_g2)
                        fig_g2.tight_layout()
                        st.pyplot(fig_g2)
                        
                        buf_g2 = io.BytesIO()
                        fig_g2.savefig(buf_g2, format="png", dpi=300, bbox_inches="tight")
                        buf_g2.seek(0)
                        st.download_button(
                            label="📸 Download Grade Chart",
                            data=buf_g2,
                            file_name="Grade_Distribution.png",
                            mime="image/png",
                            type="primary",
                            use_container_width=True,
                            key="dl_grade"
                        )
                        plt.close(fig_g2)

                st.markdown("---")
                
                # ==========================================
                # 9 & 10. Split Coil Verification 
                # ==========================================
                st.subheader("9 & 10. Split Coil Verification")
                st.info("Identifying identical coils processed on both machines to isolate machine impact.")

                coil_status_scrap = df_t6.groupby([COIL_ID_COL, 'Machine_Status']).agg({LEN_COL: 'sum', SCRAP_COL: 'sum'}).reset_index()
                coil_status_scrap['Scrap_Rate'] = np.where(coil_status_scrap[LEN_COL] > 0, (coil_status_scrap[SCRAP_COL] / coil_status_scrap[LEN_COL]) * 100, 0)
                
                old_machine_col = 'Old Machine (< Apr 2026)'
                new_machine_col = 'New Machine (>= Apr 2026)'
                
                split_pivot = coil_status_scrap.pivot(index=COIL_ID_COL, columns='Machine_Status', values='Scrap_Rate').dropna(subset=[old_machine_col, new_machine_col])
                
                if not split_pivot.empty:
                    split_pivot['Delta (%)'] = split_pivot[old_machine_col] - split_pivot[new_machine_col]
                    
                    conds = [
                        (split_pivot[old_machine_col] > 10) & (split_pivot[new_machine_col] < 5),
                        (split_pivot[old_machine_col] > 10) & (split_pivot[new_machine_col] >= 5),
                        (split_pivot[new_machine_col] > split_pivot[old_machine_col] + 5),
                        (split_pivot[old_machine_col] > 0) & (split_pivot[new_machine_col] == 0)
                    ]
                    choices = ["🚨 Old Machine Issue (Proven)", "⚠️ Material / Process Issue", "⚙️ New Machine Tuning Issue", "✅ Improved on New Machine"]
                    split_pivot['Root Cause Classification'] = np.select(conds, choices, default="✅ Normal / Stable")
                    
                    if props_cols:
                        coil_props = df_t6[df_t6[COIL_ID_COL].isin(split_pivot.index)].groupby(COIL_ID_COL)[props_cols].mean()
                        split_pivot = split_pivot.join(coil_props)
                    
                    rename_dict = {old_machine_col: 'Scrap (Old Machine)', new_machine_col: 'Scrap (New Machine)',
                                   'YS': 'Actual YS', 'TS': 'Actual TS', 'EL': 'Actual EL', 'YPE': 'Actual YPE'}
                    split_report = split_pivot.rename(columns=rename_dict).reset_index()
                    
                    format_dict = {'Scrap (Old Machine)': '{:.2f}%', 'Scrap (New Machine)': '{:.2f}%', 'Delta (%)': '{:.2f}%',
                                   'Actual YS': '{:.1f}', 'Actual TS': '{:.1f}', 'Actual EL': '{:.1f}', 'Actual YPE': '{:.1f}'}
                    st.dataframe(
                        split_report.style.format(format_dict, na_rep="N/A").background_gradient(subset=['Scrap (Old Machine)', 'Scrap (New Machine)'], cmap='Reds'),
                        use_container_width=True, hide_index=True
                    )
                else:
                    st.success("All multi-machine coils achieved perfect quality (0% scrap) or no split-coils found.")
        else:
            st.error("Missing required columns for Task 6 Analysis ('Usage Date', 'Coil ID', 'Length', or 'Scrap').")
    # ==========================================================
    # TASK 7: PRODUCTION-BASED SCRAP & MATERIAL STABILITY
    # ==========================================================
    with tab7:
        st.header("7. Production-Based Scrap & Material Stability")
        st.info("Logic: Identifies unique coils to prevent length overcounting. Length is only counted for the first occurrence of repeated coils.")

        # Tạo bản sao dữ liệu và tiền xử lý thời gian sản xuất
        df_t7 = df.dropna(subset=['Production_Date', COIL_ID_COL]).copy()
        
        # --- LỌC DỮ LIỆU TỪ QUÝ 3/2025 TRỞ ĐI ---
        df_t7 = df_t7[df_t7['Production_Date'] >= pd.Timestamp(2025, 6, 29)] # Lấy từ Q3 2025
        
        df_t7['Prod_Month'] = df_t7['Production_Date'].dt.strftime('%Y-%m')
        
        # --- XỬ LÝ DỮ LIỆU LẶP (COIL DEDUPLICATION) ---
        # Sắp xếp theo ngày sản xuất để xác định lần đầu tiên xuất hiện
        df_t7 = df_t7.sort_values([COIL_ID_COL, 'Production_Date'])
        
        # Lấy bản ghi đầu tiên của mỗi cuộn để tính Chiều dài (Input Length)
        df_unique_first = df_t7.drop_duplicates(subset=[COIL_ID_COL], keep='first')
        monthly_input_len = df_unique_first.groupby('Prod_Month')[LEN_COL].sum()
        
        # Tính tổng Scrap (Cộng dồn tất cả các lần phát sinh scrap của cuộn đó)
        monthly_total_scrap = df_t7.groupby('Prod_Month')[SCRAP_COL].sum()
        
        # Tính giá trị trung bình của cơ tính (Actual Values)
        # Sử dụng toàn bộ dữ liệu để có cái nhìn tổng quát về độ biến động
        prop_cols = [c for c in ['YS', 'TS', 'EL', 'YPE'] if c in df_t7.columns]
        monthly_props = df_t7.groupby('Prod_Month')[prop_cols].mean()
        
        # Gộp dữ liệu phân tích
        t7_summary = pd.DataFrame({
            'Input_Length': monthly_input_len,
            'Total_Scrap': monthly_total_scrap
        }).join(monthly_props).reset_index()
        
        t7_summary['Scrap_Rate (%)'] = np.where(
            t7_summary['Input_Length'] > 0,
            (t7_summary['Total_Scrap'] / t7_summary['Input_Length'] * 100),
            0
        ).round(2)

        # --- HIỂN THỊ BIỂU ĐỒ TƯƠNG QUAN ---
        st.subheader("Correlation: Scrap Rate vs. Actual Values (Factory Date)")
        
        t7_row1 = st.columns(2)
        t7_row2 = st.columns(2)
        t7_cols = t7_row1 + t7_row2
        
        features_to_plot = [
            ('YS', 'Actual YS', '#1f77b4'),
            ('TS', 'Actual TS', '#2ca02c'),
            ('EL', 'Actual EL', '#9467bd'),
            ('YPE', 'Actual YPE', '#ff7f0e')
        ]

        for idx, (feat_id, label, color) in enumerate(features_to_plot):
            if feat_id in t7_summary.columns:
                with t7_cols[idx]:
                    fig_t7, ax1 = plt.subplots(figsize=(7, 4.5))
                    
                    # Trục trái: Scrap Rate
                    ax1.set_xlabel('Production Month')
                    ax1.set_ylabel('Scrap Rate (%)', color='#d62728', fontweight='bold')
                    ax1.plot(t7_summary['Prod_Month'], t7_summary['Scrap_Rate (%)'], 
                            color='#d62728', marker='o', linewidth=2.5, label='Scrap Rate')
                    ax1.tick_params(axis='y', labelcolor='#d62728')
                    
                    # Trục phải: Actual Property
                    ax2 = ax1.twinx()
                    ax2.set_ylabel(label, color=color, fontweight='bold')
                    ax2.plot(t7_summary['Prod_Month'], t7_summary[feat_id], 
                            color=color, marker='s', linestyle='--', alpha=0.8, label=label)
                    ax2.tick_params(axis='y', labelcolor=color)
                    
                    plt.title(f"Scrap Rate vs {label}", fontweight='bold')
                    ax1.set_xticklabels(t7_summary['Prod_Month'], rotation=45, ha='right')
                    
                    if 'add_chart_border' in globals():
                        add_chart_border(ax1)
                        
                    fig_t7.tight_layout()
                    st.pyplot(fig_t7)
                    
                    # Nút tải ảnh chất lượng cao cho Task 7
                    buf_t7 = io.BytesIO()
                    fig_t7.savefig(buf_t7, format="png", dpi=300, bbox_inches="tight")
                    buf_t7.seek(0)
                    st.download_button(
                        label=f"📸 Download {label} Chart",
                        data=buf_t7,
                        file_name=f"Prod_Scrap_vs_{feat_id}.png",
                        mime="image/png",
                        type="primary",
                        use_container_width=True,
                        key=f"dl_t7_chart_{idx}" 
                    )
                    plt.close(fig_t7) # FIX: Ngăn sập RAM

        # Hiển thị bảng dữ liệu chi tiết
        st.markdown("### Production Monthly Analytics Data")
        st.dataframe(
            t7_summary.style.format({
                'Input_Length': '{:,.1f}', 'Total_Scrap': '{:,.1f}', 
                'Scrap_Rate (%)': '{:.2f}%', 'YS': '{:.1f}', 
                'TS': '{:.1f}', 'EL': '{:.2f}', 'YPE': '{:.2f}'
            }), use_container_width=True
        )
            
    # --- GLOBAL EXPORT ---
    st.sidebar.header("Export Reports")
    if st.sidebar.button("Generate Excel File"):
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            if not yield_summary.empty:
                yield_summary.to_excel(writer, sheet_name='Yield_Detailed', index=False)
            if 'grade_dist_display' in locals() and not grade_dist_display.empty:
                grade_dist_display.to_excel(writer, sheet_name='Grade_Distribution')
            if 'cap_summary_rows' in locals() and cap_summary_rows:
                pd.DataFrame(cap_summary_rows).to_excel(writer, sheet_name='Capability_Log', index=False)
            if 'plot_df_base' in locals() and not plot_df_base.empty:
                plot_df_base.to_excel(writer, sheet_name='Task4_IMR_Data', index=False)
            if 'trend_data' in locals() and not trend_data.empty:
                trend_data.drop(columns=['_sort']).to_excel(writer, sheet_name='Trend_Data', index=False)
            if 'scrap_by_period' in locals() and not scrap_by_period.empty:
                scrap_by_period.to_excel(writer, sheet_name='Scrap_By_Period', index=False)
            if 'scrap_detail' in locals() and not scrap_detail.empty:
                scrap_detail.to_excel(writer, sheet_name='Scrap_Detailed', index=False)
            if 't7_summary' in locals() and not t7_summary.empty:
                t7_summary.to_excel(writer, sheet_name='Task7_Production_Stab', index=False)
        
        st.sidebar.download_button(
            label="📥 Download Full Excel",
            data=output.getvalue(),
            file_name="Quality_Scrap_Deep_Analysis.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
